#requirements.txt

from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from langchain_community.llms import Cohere
#import FAISS
from langchain_community.vectorstores import FAISS
#from langchain_google_genai import GoogleGenerativeAI
import google.generativeai as genai
from langchain_google_genai import GoogleGenerativeAIEmbeddings

from langchain_google_genai import ChatGoogleGenerativeAI
#from langchain.embeddings import HuggingFaceEmbeddings
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser
from langchain_core.messages import AIMessage, HumanMessage
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.prompts  import PromptTemplate, ChatPromptTemplate, MessagesPlaceholder
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
import os
import time
# necessary Imports Ends
load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

pdf_file = open("D://GENAI//data//NCERT_Science.pdf",'rb')
#ppt_file = Presentation("/kaggle/input/mid-report-ppt/Nitesh_PPT.pptx")
#doc_file = Document('/kaggle/input/final-report-synopsis/final_project synopsis.docx')


#Data Extraction

# PDF :- Pdf data is extracted using PyPDF2 and all text is stored in a string.
# PPT :- PPT data is extracted using python-pptx module and all text is stored in a string.
# DOCS :- Docs data is extracted using python-docs module and all text is stored in a string.
# After Extracting all data seperately, I have combined all text in a single string for further text processing.


# extracting pdf data
pdf_text = ""
pdf_reader = PdfReader(pdf_file)
for page in pdf_reader.pages:
    pdf_text += page.extract_text()

# extracting ppt data
#ppt_text = ""
#for slide in ppt_file.slides:
#    for shape in slide.shapes:
#        if hasattr(shape, "text"):
#            ppt_text += shape.text + '\n'

## extracting doc data
#doc_text = ""
#for paragraph in doc_file.paragraphs:
#    doc_text += paragraph.text + '\n'


# merging all the text 

#all_text = pdf_text + '\n' + ppt_text + '\n' + doc_text
all_text = pdf_text 


len(all_text)



#Chunking

#In this step I am creating the chunks of data, for this step I am using Recursive Character Splitter which break large Documents into smaller chunks. 
#This is useful both for indexing data and for passing it in to a model, since large chunks are harder to search over and won’t fit in a model’s finite context window.Chunking

# splitting the text into chunks for embeddings creation

text_splitter = RecursiveCharacterTextSplitter(
        chunk_size = 1000, 
        chunk_overlap = 200, # This is helpul to handle the data loss while chunking.
        length_function = len,
        separators=['\n', '\n\n', ' ', '']
    )
    
chunks = text_splitter.split_text(text = all_text)


len(chunks)


##Embeddings Creation

##Embeddings creation is a crucial preprocessing step in the development of document-based Question and Answering (Q&A) systems. 
##This process involves converting textual data from documents and questions into dense, high-dimensional vectors known as embeddings. 
##These embeddings are designed to capture the semantic meaning of words, sentences, or even entire documents, enabling the Q&A system to understand and process natural language more effectively.

# Initializing embeddings model

embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")

#Indexing

#Indexing data using Facebook AI Similarity Search (FAISS) is a pivotal step in developing efficient and scalable document-based Question and Answering (Q&A) systems. 
# FAISS is a library that facilitates the efficient search for similarities in large datasets, especially useful for tasks involving high-dimensional vectors like text embeddings. 
# When applied to document-based Q&A, FAISS indexes the embeddings of document chunks (e.g., paragraphs, sentences) to optimize the retrieval process.


# Indexing the data using FAISS
vectorstore = FAISS.from_texts(chunks, embedding = embeddings)

#Retriever

#In the development of document-based Question and Answering (Q&A) systems, creating a retriever is a crucial step that 
# directly impacts the system's ability to find relevant information efficiently. The retriever utilizes the pre-indexed embeddings of document chunks, 
# searching through them to find the most relevant pieces of content in response to a user query. This process involves setting up a retrieval 
# mechanism that leverages similarity search to identify the best matches for the query embeddings within the indexed data.

# creating retriever
retriever = vectorstore.as_retriever(search_type="similarity", search_kwargs={"k": 6})

retrieved_docs = retriever.invoke("How did the Swadeshi Movement influence Indian industries in the early 20th century?")
len(retrieved_docs)

print(retrieved_docs[0].page_content)

#LLM Models

# Large Language Models (LLMs) are advanced artificial intelligence systems designed to understand, generate, and interact with human language in a 
# way that mimics human-like understanding. These models are trained on vast amounts of text data, allowing them to grasp the nuances of language, 
# including grammar, context, and even cultural references. The capabilities of LLMs extend beyond simple text generation; they can perform a 
# variety of tasks such as translation, summarization, question answering, and even code generation.
#One of the key technologies behind LLMs is the Transformer architecture, which enables the model to pay attention to different 
# parts of the input text differently, thereby understanding the context and relationships between words and phrases more effectively. 
# This architecture has led to significant improvements in natural language processing tasks and is the foundation of many state-of-the-art LLMs.

# Cohere LLM

prompt_template = """Answer the question as precise as possible using the provided context. If the answer is
                not contained in the context, say "answer not available in context" \n\n
                Context: \n {context}?\n
                Question: \n {question} \n
                Answer:"""

prompt = PromptTemplate.from_template(template=prompt_template)

# function to create a single string of relevant documents given by Faiss.
def format_docs(docs):
    return "\n\n".join(doc.page_content for doc in docs)

# RAG Chain

def generate_answer(question):
    model = ChatGoogleGenerativeAI(model="gemini-1.5-flash", temperature=0.1)

    rag_chain = (
        {"context": retriever | format_docs, "question": RunnablePassthrough()}
        | prompt
        | model
        | StrOutputParser()
    )
    
    return rag_chain.invoke(question)

# Results
ans = generate_answer("How did the Swadeshi Movement influence Indian industries in the early 20th century?")
print(ans)

ans = generate_answer("Who is virat kohli")
print(ans)

ans = generate_answer("How did the East India Company contribute to the opium trade with China in the 19th century?")
print(ans)

ans = generate_answer("What was the impact of British manufactured goods on the Indian market during the 19th century?")
print(ans)

ans = generate_answer("What is the primary goal of the project?")
print(ans)
      
ans = generate_answer("Which machine learning algorithms are utilized in the project?")
print(ans)

ans = generate_answer("What preprocessing techniques are used in the project?")
print(ans)

ans = generate_answer("How was the project deployed?")
print(ans)

#Conclusion

#In conclusion, this Kaggle notebook has successfully demonstrated the application of Retrieval-Augmented Generation (RAG) 
# for multi-document Question and Answering. It showcased the power of combining retrieval and generation capabilities to provide 
# accurate, context-aware answers sourced from multiple documents. Through detailed examples, performance evaluations, and 
# interactive demonstrations, the notebook highlights the efficiency and scalability of RAG in handling complex Q&A tasks.












